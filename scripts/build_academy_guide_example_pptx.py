from __future__ import annotations

from pathlib import Path


def _add_title(slide, title: str, subtitle: str = ""):
    shapes = slide.shapes
    title_box = shapes.title
    if title_box:
        title_box.text = title
    if subtitle and hasattr(slide, "placeholders") and len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = subtitle
        except Exception:
            pass


def _add_bullets(slide, heading: str, bullets: list[str], *, left=0.7, top=1.6, width=12.0, height=5.0):
    from pptx.util import Inches

    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = heading
    p.level = 0
    for b in bullets:
        pp = tf.add_paragraph()
        pp.text = b
        pp.level = 1


def build(out_path: Path):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.core_properties.title = "Ejemplo · Guía práctica (Academia)"
    prs.core_properties.subject = "Guía práctica para fútbol base"

    # Slide 1: Portada
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    _add_title(s1, "GUÍA PRÁCTICA", "Ejemplo (Academia) · Presión orientada a banda")

    # Slide 2: Estructura de la guía
    s2 = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    s2.shapes.title.text = "Cómo se ve una guía (estructura)"
    _add_bullets(
        s2,
        "Bloques típicos",
        [
            "Objetivo (en 1 frase) + por qué importa",
            "Reglas simples (3) para repetir en campo",
            "Roles: 1º, 2º, 3º defensor (o equivalentes)",
            "Triggers: cuándo sí / cuándo no",
            "Checklist rápido (para staff y jugadores)",
            "Errores típicos + corrección concreta (1 frase)",
            "Escena 2D (pizarra) para visualizar en 15s",
            "Mini-quiz (opcional) para fijar concepto",
        ],
    )

    # Slide 3: Ejemplo concreto
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    s3.shapes.title.text = "Ejemplo: Presión orientada (guiar a banda)"
    _add_bullets(
        s3,
        "Objetivo",
        [
            "Forzar al rival a banda/atrás con cobertura preparada",
            "Robar o provocar error SIN partir el equipo",
        ],
        top=1.4,
        height=1.8,
    )
    _add_bullets(
        s3,
        "3 reglas simples",
        [
            "1º defensor: carrera curva, cierra dentro, invita fuera",
            "2º defensor: tapa pase interior y achica",
            "3º defensor: equilibra espalda y cambio de orientación",
        ],
        top=3.0,
        height=3.0,
    )

    # Slide 4: Triggers + checklist
    s4 = prs.slides.add_slide(prs.slide_layouts[5])
    s4.shapes.title.text = "Triggers + checklist (para usar en partido)"
    _add_bullets(
        s4,
        "Triggers para SALTAR",
        [
            "Receptor de espaldas",
            "Pase lento o flotado",
            "Control malo / balón se separa del pie",
            "Interior cerrado (no hay pase por dentro)",
        ],
        top=1.4,
        height=2.6,
    )
    _add_bullets(
        s4,
        "Checklist (10 segundos)",
        [
            "¿Hay cobertura? (2º defensor listo)",
            "¿Está cerrada la línea interior?",
            "¿Si me superan, quién protege la espalda?",
            "Si NO: temporiza y orienta sin morder",
        ],
        top=3.9,
        height=2.6,
    )

    # Slide 5: Errores típicos + frases de coaching
    s5 = prs.slides.add_slide(prs.slide_layouts[5])
    s5.shapes.title.text = "Errores típicos (y frase corta de corrección)"
    _add_bullets(
        s5,
        "Errores",
        [
            "1º defensor se lanza sin ayudas → 'TEMPORIZA, CERRANDO DENTRO'",
            "Bloque no acompaña → 'ACORTA Y ACOMPAÑA'",
            "Se abre el pasillo interior → 'DENTRO CERRADO'",
            "Nadie al rechace / segunda jugada → 'RECHACE ES NUESTRO'",
        ],
        top=1.4,
        height=4.6,
    )

    # Estilo mínimo: color en títulos
    for slide in prs.slides:
        try:
            if slide.shapes.title and slide.shapes.title.text_frame:
                for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(15, 23, 42)  # slate-900
                    run.font.size = Pt(34)
        except Exception:
            pass

    # Ajuste portada
    try:
        t = s1.shapes.title
        if t and t.text_frame:
            t.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            for run in t.text_frame.paragraphs[0].runs:
                run.font.size = Pt(44)
                run.font.bold = True
                run.font.color.rgb = RGBColor(2, 132, 199)  # sky-600
    except Exception:
        pass

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


if __name__ == "__main__":
    here = Path(__file__).resolve()
    repo = here.parents[1]
    out = repo / "artifacts" / "ejemplo_guia_practica_presion_orientada.pptx"
    build(out)
    print(str(out))

